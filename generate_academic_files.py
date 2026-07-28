from pathlib import Path
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.shared import Inches, Pt, RGBColor
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as PRGB
from pptx.util import Inches as PInches, Pt as PPt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
BLUE = "2563EB"; NAVY = "0F172A"; SLATE = "475569"; LIGHT = "EFF6FF"


def add_page_number(paragraph):
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run()
    fld = OxmlElement("w:fldSimple"); fld.set(qn("w:instr"), "PAGE")
    run._r.append(fld)


def build_docx():
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(.75); sec.bottom_margin = Inches(.7)
    styles = doc.styles
    styles["Normal"].font.name = "Aptos"; styles["Normal"].font.size = Pt(10.5)
    for name, size, color in [("Title", 28, NAVY), ("Heading 1", 19, BLUE), ("Heading 2", 14, NAVY)]:
        styles[name].font.name = "Aptos Display"; styles[name].font.size = Pt(size); styles[name].font.color.rgb = RGBColor.from_string(color)
    footer = sec.footer.paragraphs[0]
    footer.add_run("AI Resume Analyzer  •  Academic Project Report  •  ")
    add_page_number(footer)

    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("AI RESUME ANALYZER").bold = True; p.runs[0].font.size = Pt(30); p.runs[0].font.color.rgb = RGBColor.from_string(BLUE)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Academic Project Report"); r.font.size = Pt(20); r.font.color.rgb = RGBColor.from_string(NAVY)
    doc.add_paragraph("\n")
    table = doc.add_table(rows=5, cols=2); table.style = "Light Shading Accent 1"
    fields = [("Project", "AI Resume Analyzer"), ("Technology", "Python, Streamlit, spaCy, Hugging Face"), ("Inputs", "Resume file + pasted/uploaded job description"), ("Outputs", "ATS, keywords, skills, semantic alignment"), ("Date", "July 2026")]
    for row, (a,b) in zip(table.rows, fields): row.cells[0].text=a; row.cells[1].text=b
    doc.add_page_break()

    abstract = (DOCS / "ABSTRACT.md").read_text(encoding="utf-8").splitlines()
    report = (DOCS / "ACADEMIC_REPORT.md").read_text(encoding="utf-8").splitlines()
    lines = abstract + ["", "# Table of Contents", ""] + report[1:]
    for line in lines:
        s=line.strip()
        if not s: continue
        if s.startswith("# "):
            doc.add_heading(s[2:], level=1)
        elif s.startswith("## "):
            doc.add_heading(s[3:], level=2)
        elif s.startswith("- "):
            doc.add_paragraph(s[2:], style="List Bullet")
        elif s.startswith("**Keywords:**"):
            p=doc.add_paragraph(); p.add_run("Keywords: ").bold=True; p.add_run(s.replace("**Keywords:**", ""))
        else:
            doc.add_paragraph(s)
    # Add an architecture table before conclusion area as appendix.
    doc.add_heading("Appendix A: Architecture Summary", level=1)
    t=doc.add_table(rows=1, cols=3); t.style="Light Shading Accent 1"
    for c, text in zip(t.rows[0].cells, ["Layer", "Technology", "Purpose"]): c.text=text
    for values in [("UI","Streamlit","Upload, configure, visualize, export"),("Parsing","PyPDF2 / python-docx / antiword","Extract text"),("NLP","spaCy / taxonomy","Keywords and skills"),("Semantic","Hugging Face MiniLM","Contextual similarity"),("Scoring","Python rules","Explainable scores and actions")]:
        cells=t.add_row().cells
        for c,v in zip(cells,values): c.text=v
    path=DOCS/"AI_Resume_Analyzer_Academic_Report.docx"; doc.save(path); return path


def add_slide(prs, title, bullets, number, accent=PRGB(37,99,235)):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=PRGB(248,250,252)
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,PInches(.22),prs.slide_height); bar.fill.solid(); bar.fill.fore_color.rgb=accent; bar.line.fill.background()
    tx=slide.shapes.add_textbox(PInches(.7),PInches(.48),PInches(11.8),PInches(.8)).text_frame
    p=tx.paragraphs[0]; p.text=title; p.font.size=PPt(28); p.font.bold=True; p.font.color.rgb=PRGB(15,23,42)
    body=slide.shapes.add_textbox(PInches(.9),PInches(1.55),PInches(11.3),PInches(5.1)).text_frame
    body.word_wrap=True
    for i,item in enumerate(bullets):
        p=body.paragraphs[0] if i==0 else body.add_paragraph(); p.text=item; p.font.size=PPt(19); p.font.color.rgb=PRGB(51,65,85); p.space_after=PPt(15); p.level=0
    foot=slide.shapes.add_textbox(PInches(11.7),PInches(7),PInches(.7),PInches(.3)).text_frame.paragraphs[0]
    foot.text=str(number); foot.font.size=PPt(10); foot.font.color.rgb=PRGB(100,116,139); foot.alignment=PP_ALIGN.RIGHT
    return slide


def build_pptx():
    prs=Presentation(); prs.slide_width=PInches(13.333); prs.slide_height=PInches(7.5)
    slide=prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb=PRGB(15,23,42)
    shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PInches(.85),PInches(1.0),PInches(1.25),PInches(.14)); shape.fill.solid(); shape.fill.fore_color.rgb=PRGB(59,130,246); shape.line.fill.background()
    tb=slide.shapes.add_textbox(PInches(.85),PInches(1.35),PInches(11.4),PInches(2)).text_frame
    p=tb.paragraphs[0]; p.text="AI Resume Analyzer"; p.font.size=PPt(40); p.font.bold=True; p.font.color.rgb=PRGB(255,255,255)
    p=tb.add_paragraph(); p.text="Explainable ATS, keyword, skill, and semantic alignment"; p.font.size=PPt(22); p.font.color.rgb=PRGB(147,197,253); p.space_before=PPt(14)
    p=tb.add_paragraph(); p.text="Python • Streamlit • spaCy • Hugging Face Transformers"; p.font.size=PPt(16); p.font.color.rgb=PRGB(203,213,225); p.space_before=PPt(30)

    slides=[
      ("Problem and motivation",["Applicants cannot easily verify resume parseability or vacancy alignment.","Proprietary ATS behavior varies; a universal pass/fail claim would be misleading.","Need: transparent evidence, practical gaps, and ethical improvement guidance."]),
      ("Project objectives",["Accept PDF/DOCX/DOC resumes and pasted or uploaded job descriptions.","Evaluate ATS-oriented structure and extract job terminology.","Combine exact skill matching with contextual transformer similarity.","Return explainable scores, actions, and downloadable results."]),
      ("User workflow",["1  Upload resume","2  Paste JD, upload JD, or combine both","3  Select transformer or lightweight mode","4  Analyze and inspect component evidence","5  Download TXT or JSON report"]),
      ("System architecture",["Presentation: Streamlit interface and validation","Ingestion: PyPDF2, python-docx, antiword","NLP: normalization, spaCy, skill taxonomy and aliases","Semantics: Hugging Face MiniLM with TF-IDF fallback","Output: weighted scores, findings, tables, recommendations"]),
      ("NLP and matching pipeline",["Normalize case, whitespace, punctuation, and technical symbols.","Boundary-aware multiword phrase matching avoids substring errors.","Canonical aliases map AWS, PostgreSQL, React.js, NLP, and more.","Rank job unigrams/bigrams; compare resume presence.","Mean-pool MiniLM token vectors and compute cosine similarity."]),
      ("Explainable scoring",["ATS compatibility — 25%","Keyword coverage — 30%","Weighted skill alignment — 30%","Semantic similarity — 15%","Mandatory / normal / preferred skill weights: 1.5 / 1.0 / 0.7"]),
      ("ATS compatibility checks",["Parseability and extracted text length (20)","Email and phone in plain text (20)","Standard summary, experience, education, skills, projects headings (25)","Formatting safety and file type (20)","Action verbs and quantified achievements (15)"]),
      ("Technology stack",["Python 3.11; modular, testable implementation","Streamlit dashboard; Pandas tables and charts","PyPDF2 and python-docx; antiword for legacy DOC","spaCy plus configurable JSON skill catalog","Transformers + PyTorch MiniLM; scikit-learn fallback"]),
      ("Testing and evaluation",["Unit tests: parsing, DOCX tables, aliases, score ranges, gaps.","Integration cases: corrupt/encrypted/scanned files and complex layouts.","Future labeled benchmark: precision, recall, F1, human-score correlation.","Measure cold start, inference latency, failure rate, and subgroup errors."]),
      ("Privacy, security, and ethics",["Reference app processes in memory and deliberately stores no documents.","Production: scan files, verify type, sandbox parsers, encrypt and expire data.","Do not infer protected traits or use the score as an automated hiring gate.","Recommend only truthful skills; disclose that scores are guidance."]),
      ("Limitations and future work",["Scanned PDFs need OCR; graphic layouts may extract out of order.","Finite taxonomy and 512-token transformer context constrain coverage.","Next: section-aware chunk embeddings, multilingual support, ESCO/O*NET.","Prioritize consented evaluation and fairness governance over model size."]),
      ("Conclusion",["A working end-to-end academic prototype with both JD input modes.","Combines rule-based transparency, lexical coverage, and semantic NLP.","Modular architecture supports deployment and responsible extension.","Deliverables: source, tests, README, plan, report, specification, slides."]),
    ]
    for i,(title,bullets) in enumerate(slides,2): add_slide(prs,title,bullets,i)
    path=DOCS/"AI_Resume_Analyzer_Presentation.pptx"; prs.save(path); return path

if __name__=="__main__":
    print(build_docx())
    print(build_pptx())
